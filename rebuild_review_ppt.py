#!/usr/bin/env python3
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

LOGO = "/tmp/ppt_build/media/sheridan_logo.png"
OUT = "/Users/lulu/Downloads/Business_Statistics_Post_Midterm_Review.pptx"

DARK = RGBColor(0x1E, 0x4E, 0x79)
FILL = RGBColor(0xF5, 0xF9, 0xFF)
BORDER = RGBColor(0xAA, 0xC8, 0xEB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def style_run(run, size=18, bold=False):
    run.font.name = "Arial"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = DARK


def add_logo(slide):
    slide.shapes.add_picture(LOGO, Inches(8.3), Inches(0.15), height=Inches(0.55))


def add_title(slide, text):
    box = slide.shapes.add_textbox(Inches(0.45), Inches(0.2), Inches(7.5), Inches(0.6))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    style_run(run, size=28, bold=True)


def rounded_box(slide, left, top, width, height):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = FILL
    shape.line.color.rgb = BORDER
    shape.line.width = Pt(1.5)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.15)
    tf.margin_bottom = Inches(0.15)
    return tf


def fill_box(tf, lines):
    tf.clear()
    for i, item in enumerate(lines):
        if isinstance(item, tuple):
            text, level, bold, size = item
        else:
            text, level, bold, size = item, 0, False, 16
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = level
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = text
        style_run(run, size=size, bold=bold)


def blank_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    add_logo(slide)
    return slide


def chapter_slide(prs, title, left_lines, right_lines):
    slide = blank_slide(prs)
    add_title(slide, title)
    fill_box(rounded_box(slide, Inches(0.5), Inches(1.0), Inches(4.5), Inches(5.0)), left_lines)
    fill_box(rounded_box(slide, Inches(5.3), Inches(1.0), Inches(4.5), Inches(5.0)), right_lines)


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    s1 = blank_slide(prs)
    add_title(s1, "📊 Business Statistics — Post-Midterm Review")
    fill_box(
        rounded_box(s1, Inches(0.55), Inches(1.0), Inches(8.9), Inches(4.5)),
        [
            ("Tutorial Review Pack", 0, True, 20),
            ("✓ Chapters 1–7 Recap", 1, False, 16),
            ("✓ Key Concepts & Formulas", 1, False, 16),
            ("✓ Worked Examples", 1, False, 16),
            ("✓ Excel Cheat Sheet", 1, False, 16),
            ("✓ Common Mistakes", 1, False, 16),
        ],
    )

    chapters = [
        (
            "Chapter 1 — Introduction to Statistics",
            [
                ("Key Concepts", 0, True, 18),
                ("Population vs Sample", 1, False, 15),
                ("Parameter vs Statistic", 1, False, 15),
                ("Descriptive vs Inferential", 1, False, 15),
                ("Qualitative vs Quantitative", 1, False, 15),
                ("Levels of Measurement", 1, False, 15),
            ],
            [
                ("Worked Example", 0, True, 18),
                ("Population: All Sheridan students", 1, False, 15),
                ("Sample: 120 surveyed students", 1, False, 15),
            ],
        ),
        (
            "Chapter 2 — Tables & Graphs",
            [
                ("Key Concepts", 0, True, 18),
                ("Frequency Distribution", 1, False, 15),
                ("Histogram & Bar Chart", 1, False, 15),
                ("Stem Plot & Pie Chart", 1, False, 15),
                ("Scatter Plot & Correlation", 1, False, 15),
            ],
            [
                ("Tip", 0, True, 18),
                ("Histogram → Continuous Data", 1, False, 15),
                ("Bar Chart → Categorical Data", 1, False, 15),
                ("Data Analysis → Histogram", 1, False, 15),
            ],
        ),
        (
            "Chapter 3 — Descriptive Statistics",
            [
                ("Key Concepts", 0, True, 18),
                ("Mean, Median, Mode", 1, False, 15),
                ("Variance & Std Dev", 1, False, 15),
                ("Quartiles & IQR", 1, False, 15),
                ("Outliers & Boxplots", 1, False, 15),
            ],
            [
                ("Worked Example", 0, True, 18),
                ("Q1=20  Q3=40", 1, False, 15),
                ("IQR=20", 1, False, 15),
                ("Upper Fence=70", 1, False, 15),
            ],
        ),
        (
            "Chapter 4 — Probability",
            [
                ("Key Concepts", 0, True, 18),
                ("Complement Rule", 1, False, 15),
                ("Addition Rule", 1, False, 15),
                ("Multiplication Rule", 1, False, 15),
                ("Conditional Probability", 1, False, 15),
                ("Independent Events", 1, False, 15),
            ],
            [
                ("Worked Example", 0, True, 18),
                ("P(A)=0.6  P(B)=0.5", 1, False, 15),
                ("Independent", 1, False, 15),
                ("P(A∩B)=0.30", 1, False, 15),
            ],
        ),
        (
            "Chapter 5 — Discrete Distributions",
            [
                ("Key Concepts", 0, True, 18),
                ("Binomial", 1, False, 15),
                ("Poisson", 1, False, 15),
                ("Geometric", 1, False, 15),
                ("Hypergeometric", 1, False, 15),
            ],
            [
                ("Binomial Conditions", 0, True, 18),
                ("Fixed n", 1, False, 15),
                ("Independent trials", 1, False, 15),
                ("Same probability p", 1, False, 15),
            ],
        ),
        (
            "Chapter 6 — Normal Distribution",
            [
                ("Key Concepts", 0, True, 18),
                ("Uniform Distribution", 1, False, 15),
                ("Z-score", 1, False, 15),
                ("Left/Right Tail", 1, False, 15),
                ("Between & Inverse", 1, False, 15),
                ("Percentile", 1, False, 15),
            ],
            [
                ("Worked Example", 0, True, 18),
                ("P(85<X<120)", 1, False, 15),
                ("=NORM.DIST(120,100,15,TRUE)", 1, False, 14),
                ("-NORM.DIST(85,100,15,TRUE)=0.7495", 1, False, 14),
            ],
        ),
        (
            "Chapter 7 — Confidence Intervals",
            [
                ("Key Concepts", 0, True, 18),
                ("CI for Mean", 1, False, 15),
                ("CI for Proportion", 1, False, 15),
                ("Sample Size", 1, False, 15),
                ("90%→1.645  95%→1.96", 1, False, 15),
            ],
            [
                ("Formula", 0, True, 18),
                ("CI = x̄ ± z·σ/√n", 1, False, 15),
                ("=NORM.S.INV for z-value", 1, False, 15),
            ],
        ),
    ]

    for title, left, right in chapters:
        chapter_slide(prs, title, left, right)

    s9 = blank_slide(prs)
    add_title(s9, "💻 Excel Cheat Sheet")
    fill_box(
        rounded_box(s9, Inches(0.5), Inches(1.0), Inches(4.5), Inches(4.8)),
        [
            ("Probability", 0, True, 18),
            ("=NORM.DIST(x,mean,sd,TRUE)", 1, False, 14),
            ("=NORM.S.DIST(z,TRUE)", 1, False, 14),
            ("=NORM.INV(p,mean,sd)", 1, False, 14),
            ("=NORM.S.INV(p)", 1, False, 14),
        ],
    )
    fill_box(
        rounded_box(s9, Inches(5.3), Inches(1.0), Inches(4.5), Inches(4.8)),
        [
            ("Statistics", 0, True, 18),
            ("=AVERAGE()  =MEDIAN()", 1, False, 14),
            ("=STDEV.S()", 1, False, 14),
            ("=PERCENTILE.INC()", 1, False, 14),
            ("=QUARTILE.INC()", 1, False, 14),
            ("=CONFIDENCE.NORM()", 1, False, 14),
        ],
    )

    s10 = blank_slide(prs)
    add_title(s10, "⚠ Common Midterm Mistakes")
    fill_box(
        rounded_box(s10, Inches(0.5), Inches(1.0), Inches(8.9), Inches(4.8)),
        [
            ("Don't Lose Easy Marks", 0, True, 18),
            ("• DIST vs INV", 1, False, 15),
            ("• Left tail vs Right tail", 1, False, 15),
            ("• PERCENTILE.INC vs EXC", 1, False, 15),
            ("• Percentile ≠ Probability", 1, False, 15),
            ("• Don't forget Complement Rule", 1, False, 15),
            ("• Wrong distribution (Binomial vs Normal)", 1, False, 15),
            ("• Convert to Z-score first", 1, False, 15),
        ],
    )

    s11 = blank_slide(prs)
    add_title(s11, "📢 Tutorial Update — Going Forward")
    fill_box(
        rounded_box(s11, Inches(0.5), Inches(1.0), Inches(8.9), Inches(5.0)),
        [
            ("Tutorial Videos on Slate", 0, True, 18),
            ("• Step-by-step videos posted on Slate", 1, False, 15),
            ("• Watch videos before coming to tutorial", 1, False, 15),
            ("• Covers Ch 1–7 examples & Excel tips", 1, False, 15),
            ("In-Person Tutorial = Q&A Only", 0, True, 18),
            ("• Offline time for Questions & Answers", 1, False, 15),
            ("• Bring specific questions from videos/homework", 1, False, 15),
            ("• We clarify concepts — not re-teach the lesson", 1, False, 15),
        ],
    )

    prs.save(OUT)
    print(f"Rebuilt: {OUT} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
