#!/usr/bin/env python3
import shutil
from copy import deepcopy

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree

DARK = "1E4E79"
SRC = "/Users/lulu/Downloads/Business_Statistics_Midterm_Tutorial_Pack_with_Sheridan_Logo.pptx"
OUT = "/Users/lulu/Downloads/Business_Statistics_Post_Midterm_Review.pptx"


def set_para_color(paragraph, hex_color=DARK, bold=None, size=None):
    pPr = paragraph._p.get_or_add_pPr()
    defRPr = pPr.find(qn("a:defRPr"))
    if defRPr is None:
        defRPr = etree.SubElement(pPr, qn("a:defRPr"))
    if bold is not None:
        defRPr.set("b", "1" if bold else "0")
    if size is not None:
        defRPr.set("sz", str(size))
    solidFill = defRPr.find(qn("a:solidFill"))
    if solidFill is None:
        solidFill = etree.SubElement(defRPr, qn("a:solidFill"))
    else:
        for child in list(solidFill):
            solidFill.remove(child)
    srgb = etree.SubElement(solidFill, qn("a:srgbClr"))
    srgb.set("val", hex_color)


def write_lines(shape, lines, heading_bold=True):
    tf = shape.text_frame
    existing = list(tf.paragraphs)
    for i, item in enumerate(lines):
        if isinstance(item, tuple):
            text, level, bold, size = item
        else:
            text = item
            level = 0 if i == 0 else 1
            bold = i == 0 and heading_bold
            size = None
        p = existing[i] if i < len(existing) else tf.add_paragraph()
        p.level = level
        if p.runs:
            p.runs[0].text = text
            for run in p.runs[1:]:
                run.text = ""
        else:
            p.text = text
        set_para_color(p, DARK, bold=bold, size=size)
        if p.runs:
            p.runs[0].font.color.rgb = RGBColor.from_string(DARK)
    for j in range(len(lines), len(tf.paragraphs)):
        p = tf.paragraphs[j]
        if p.runs:
            p.runs[0].text = ""
        else:
            p.text = ""


def main():
    shutil.copy(SRC, OUT)
    prs = Presentation(OUT)

    write_lines(
        prs.slides[0].shapes[0],
        [("📊 Business Statistics — Post-Midterm Review", 0, True, 2800)],
    )
    write_lines(
        prs.slides[0].shapes[1],
        [
            ("Tutorial Review Pack", 0, True, 1800),
            ("✓ Chapters 1–7 Recap", 1, False, 1500),
            ("✓ Key Concepts & Formulas", 1, False, 1500),
            ("✓ Worked Examples", 1, False, 1500),
            ("✓ Excel Cheat Sheet", 1, False, 1500),
            ("✓ Common Mistakes", 1, False, 1500),
        ],
    )

    chapters = [
        (
            "Chapter 1 — Introduction to Statistics",
            [
                ("Key Concepts", 0, True, None),
                ("Population vs Sample", 1, False, None),
                ("Parameter vs Statistic", 1, False, None),
                ("Descriptive vs Inferential", 1, False, None),
                ("Qualitative vs Quantitative", 1, False, None),
                ("Levels of Measurement", 1, False, None),
            ],
            [
                ("Worked Example", 0, True, None),
                ("Population: All Sheridan students", 1, False, None),
                ("Sample: 120 surveyed students", 1, False, None),
            ],
        ),
        (
            "Chapter 2 — Tables & Graphs",
            [
                ("Key Concepts", 0, True, None),
                ("Frequency Distribution", 1, False, None),
                ("Histogram & Bar Chart", 1, False, None),
                ("Stem Plot & Pie Chart", 1, False, None),
                ("Scatter Plot & Correlation", 1, False, None),
            ],
            [
                ("Tip", 0, True, None),
                ("Histogram → Continuous Data", 1, False, None),
                ("Bar Chart → Categorical Data", 1, False, None),
                ("Data Analysis → Histogram", 1, False, None),
            ],
        ),
        (
            "Chapter 3 — Descriptive Statistics",
            [
                ("Key Concepts", 0, True, None),
                ("Mean, Median, Mode", 1, False, None),
                ("Variance & Std Dev", 1, False, None),
                ("Quartiles & IQR", 1, False, None),
                ("Outliers & Boxplots", 1, False, None),
            ],
            [
                ("Worked Example", 0, True, None),
                ("Q1=20  Q3=40", 1, False, None),
                ("IQR=20", 1, False, None),
                ("Upper Fence=70", 1, False, None),
            ],
        ),
        (
            "Chapter 4 — Probability",
            [
                ("Key Concepts", 0, True, None),
                ("Complement Rule", 1, False, None),
                ("Addition Rule", 1, False, None),
                ("Multiplication Rule", 1, False, None),
                ("Conditional Probability", 1, False, None),
                ("Independent Events", 1, False, None),
            ],
            [
                ("Worked Example", 0, True, None),
                ("P(A)=0.6  P(B)=0.5", 1, False, None),
                ("Independent", 1, False, None),
                ("P(A∩B)=0.30", 1, False, None),
            ],
        ),
        (
            "Chapter 5 — Discrete Distributions",
            [
                ("Key Concepts", 0, True, None),
                ("Binomial", 1, False, None),
                ("Poisson", 1, False, None),
                ("Geometric", 1, False, None),
                ("Hypergeometric", 1, False, None),
            ],
            [
                ("Binomial Conditions", 0, True, None),
                ("Fixed n", 1, False, None),
                ("Independent trials", 1, False, None),
                ("Same probability p", 1, False, None),
            ],
        ),
        (
            "Chapter 6 — Normal Distribution",
            [
                ("Key Concepts", 0, True, None),
                ("Uniform Distribution", 1, False, None),
                ("Z-score", 1, False, None),
                ("Left/Right Tail", 1, False, None),
                ("Between & Inverse", 1, False, None),
                ("Percentile", 1, False, None),
            ],
            [
                ("Worked Example", 0, True, None),
                ("P(85<X<120)", 1, False, None),
                ("=NORM.DIST(120,100,15,TRUE)", 1, False, None),
                ("-NORM.DIST(85,100,15,TRUE)=0.7495", 1, False, None),
            ],
        ),
        (
            "Chapter 7 — Confidence Intervals",
            [
                ("Key Concepts", 0, True, None),
                ("CI for Mean", 1, False, None),
                ("CI for Proportion", 1, False, None),
                ("Sample Size", 1, False, None),
                ("90%→1.645  95%→1.96", 1, False, None),
            ],
            [
                ("Formula", 0, True, None),
                ("CI = x̄ ± z·σ/√n", 1, False, None),
                ("=NORM.S.INV for z-value", 1, False, None),
            ],
        ),
    ]

    for idx, (title, left, right) in enumerate(chapters):
        slide = prs.slides[idx + 1]
        write_lines(slide.shapes[0], [(title, 0, True, 2800)])
        write_lines(slide.shapes[1], left)
        write_lines(slide.shapes[2], right)

    slide9 = prs.slides[8]
    write_lines(slide9.shapes[0], [("💻 Excel Cheat Sheet", 0, True, 2800)])
    write_lines(
        slide9.shapes[1],
        [
            ("Probability", 0, True, None),
            ("=NORM.DIST(x,mean,sd,TRUE)", 1, False, None),
            ("=NORM.S.DIST(z,TRUE)", 1, False, None),
            ("=NORM.INV(p,mean,sd)", 1, False, None),
            ("=NORM.S.INV(p)", 1, False, None),
        ],
    )
    write_lines(
        slide9.shapes[2],
        [
            ("Statistics", 0, True, None),
            ("=AVERAGE()  =MEDIAN()", 1, False, None),
            ("=STDEV.S()", 1, False, None),
            ("=PERCENTILE.INC()", 1, False, None),
            ("=QUARTILE.INC()", 1, False, None),
            ("=CONFIDENCE.NORM()", 1, False, None),
        ],
    )

    slide10 = prs.slides[9]
    write_lines(slide10.shapes[0], [("⚠ Common Midterm Mistakes", 0, True, 2800)])
    write_lines(
        slide10.shapes[1],
        [
            ("Don't Lose Easy Marks", 0, True, None),
            ("• DIST vs INV", 1, False, None),
            ("• Left tail vs Right tail", 1, False, None),
            ("• PERCENTILE.INC vs EXC", 1, False, None),
            ("• Percentile ≠ Probability", 1, False, None),
            ("• Don't forget Complement Rule", 1, False, None),
            ("• Wrong distribution (Binomial vs Normal)", 1, False, None),
            ("• Convert to Z-score first", 1, False, None),
        ],
    )

    slide10_el = slide10._element
    s11 = prs.slides.add_slide(prs.slide_layouts[6])
    sp_tree = s11._element.find(qn("p:cSld")).find(qn("p:spTree"))
    for sp in slide10_el.findall(".//" + qn("p:sp")):
        sp_tree.append(deepcopy(sp))
    for pic in slide10_el.findall(".//" + qn("p:pic")):
        sp_tree.append(deepcopy(pic))

    text_shapes = [sh for sh in s11.shapes if sh.has_text_frame]
    write_lines(text_shapes[0], [("📢 Tutorial Update — Going Forward", 0, True, 2800)])
    write_lines(
        text_shapes[1],
        [
            ("Tutorial Videos on Slate", 0, True, None),
            ("• Step-by-step videos posted on Slate", 1, False, None),
            ("• Watch videos before coming to tutorial", 1, False, None),
            ("• Covers Ch 1–7 examples & Excel tips", 1, False, None),
            ("In-Person Tutorial = Q&A Only", 0, True, None),
            ("• Offline time for Questions & Answers", 1, False, None),
            ("• Bring specific questions from videos/homework", 1, False, None),
            ("• We clarify concepts — not re-teach the lesson", 1, False, None),
        ],
    )

    prs.save(OUT)
    print(f"Saved {OUT} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
